"""
Assemble docs/deck/slides/*.png into docs/PlayHack.pptx.

Each slide is placed full-bleed on a 16:9 canvas. The slides are rendered from
HTML rather than laid out with shapes so the typography, the palette and the
figures match the product and the README exactly — the trade is that text in
the .pptx is not editable, which is the right trade for a submission deck.
"""
import glob
import os
from pptx import Presentation
from pptx.util import Emu

W, H = Emu(12192000), Emu(6858000)          # 13.333in x 7.5in — 16:9
SLIDES = sorted(glob.glob("docs/deck/slides/slide-*.png"))
OUT = "docs/PlayHack.pptx"

if not SLIDES:
    raise SystemExit("no slides rendered — run scripts/deck.mts first")

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]

for png in SLIDES:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png, 0, 0, width=W, height=H)

prs.save(OUT)
print(f"{OUT}  ·  {len(SLIDES)} slides  ·  {os.path.getsize(OUT) // 1024} KB")
