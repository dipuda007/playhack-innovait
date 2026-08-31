"""
Build docs/PlayHack.pptx with NATIVE PowerPoint text.

Why this exists
---------------
The first deck placed a rendered PNG on every slide. That gives pixel-exact
typography, and it was the wrong trade: an image resamples whenever the slide
is scaled, so "fit to screen" always looks soft however many pixels are in it.
Native text is vector — sharp at any zoom, on any projector, and editable.

Type sizes here are set in POINTS against a 13.33 x 7.5in slide, not in pixels
against a 1600px canvas. Body copy is 16-18pt; nothing that matters is below
12pt. The previous deck had 20px body, which lands near 9pt on a real slide.

Fonts are deliberately Georgia and Arial rather than Playfair and Inter. Those
are installed on every Windows and macOS machine, so the deck cannot fall back
to something unintended on a judge's laptop. The photographic and diagram
slides still carry the real typefaces, because those are images.

  python scripts/deck_native.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

W, H = Inches(13.333), Inches(7.5)
INK   = RGBColor(0x22, 0x20, 0x1E)
INK2  = RGBColor(0x55, 0x50, 0x4A)
INK3  = RGBColor(0x83, 0x7C, 0x73)
PAPER = RGBColor(0xF7, 0xF4, 0xEE)
WASH  = RGBColor(0xF2, 0xEE, 0xE6)
RULE  = RGBColor(0xE2, 0xDC, 0xD1)
RUST  = RGBColor(0xB4, 0x51, 0x2A)
BRASS = RGBColor(0xC2, 0x9A, 0x4B)
OPEN  = RGBColor(0x2E, 0x6B, 0x3E)
TAKEN = RGBColor(0x8C, 0x2F, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF, SANS, MONO = "Georgia", "Arial", "Consolas"
M = Inches(0.72)                      # page margin
CW = W - 2 * M                        # content width

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, W, H)          # 1 = rectangle
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, size=16, color=INK2, font=SANS, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.25, caps=False, space_after=6):
    """runs: a string, or a list of (text, {overrides}) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) else [runs]
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        chunks = para if isinstance(para, list) else [(para, {})]
        for body, over in chunks:
            r = p.add_run()
            r.text = body.upper() if caps else body
            f = r.font
            f.size = Pt(over.get("size", size))
            f.bold = over.get("bold", bold)
            f.name = over.get("font", font)
            f.color.rgb = over.get("color", color)
    return tb


def rule(s, x, y, w=Inches(0.9), h=Pt(3), color=BRASS):
    r = s.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def head(s, kicker, title, n=None, foot=None):
    text(s, M, Inches(0.52), CW, Inches(0.3), kicker, size=12, color=RUST,
         bold=True, caps=True)
    text(s, M, Inches(0.86), CW, Inches(1.0), title, size=34, color=INK,
         font=SERIF, bold=True, caps=True, spacing=1.02)
    rule(s, M, Inches(1.72))
    if foot or n:
        text(s, M, H - Inches(0.62), CW - Inches(0.6), Inches(0.3),
             foot or "PlayHack · Team InnovAIT", size=10, color=INK3, caps=True)
    if n:
        text(s, W - M - Inches(0.6), H - Inches(0.62), Inches(0.6), Inches(0.3),
             n, size=10, color=INK3, align=PP_ALIGN.RIGHT)


def table(s, x, y, w, rows, widths, size=13, header=True, row_h=Inches(0.42)):
    nr, nc = len(rows), len(rows[0])
    shape = s.shapes.add_table(nr, nc, x, y, w, row_h * nr)
    tbl = shape.table
    total = sum(widths)
    for i, ww in enumerate(widths):
        tbl.columns[i].width = Emu(int(w * ww / total))
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = row_h
        for ci, cell in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid(); c.fill.fore_color.rgb = PAPER
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.04); c.margin_bottom = Inches(0.04)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.line_spacing = 1.18
            body, over = (cell, {}) if isinstance(cell, str) else cell
            r = p.add_run(); r.text = body
            f = r.font
            hdr = header and ri == 0
            f.size = Pt(over.get("size", 11 if hdr else size))
            f.bold = over.get("bold", hdr)
            f.name = over.get("font", SANS)
            f.color.rgb = over.get("color", INK3 if hdr else INK2)
            if over.get("align") == "r":
                p.alignment = PP_ALIGN.RIGHT
    return shape


def code(s, x, y, w, h, lines, size=13):
    box = s.shapes.add_shape(1, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = WASH
    box.line.color.rgb = INK; box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = tf.margin_bottom = Inches(0.16)
    tf.word_wrap = False
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.4
        for body, over in (line if isinstance(line, list) else [(line, {})]):
            r = p.add_run(); r.text = body
            f = r.font; f.size = Pt(size); f.name = MONO
            f.color.rgb = over.get("color", INK2); f.bold = over.get("bold", False)
    return box


def stat(s, x, y, value, label, color=INK, vsize=52, w=Inches(3.4)):
    text(s, x, y, w, Inches(0.8), value, size=vsize, color=color, font=MONO, bold=True)
    text(s, x, y + Inches(0.78), w, Inches(0.9), label, size=14, color=INK2, spacing=1.2)


# ── 01 · title ────────────────────────────────────────────────────────────
s = slide(INK)
s.shapes.add_picture("docs/media/deck-title-bg.png", 0, 0, width=W, height=H)
text(s, M, Inches(4.05), CW, Inches(0.3), "PlayHack · SDE Track · IIT Guwahati",
     size=13, color=RGBColor(0xE0, 0xA2, 0x67), bold=True, caps=True)
text(s, M, Inches(4.45), CW, Inches(1.1), "PLAYHACK", size=54, color=PAPER,
     font=SERIF, bold=True)
text(s, M, Inches(5.55), Inches(8.2), Inches(1.0),
     "Campus sports facility booking where the winner of a race for one court "
     "is decided by Postgres, not by our code.",
     size=18, color=RGBColor(0xE8, 0xE3, 0xDA), spacing=1.35)
text(s, M, H - Inches(0.78), CW, Inches(0.3),
     "Team InnovAIT     ·     innovait-hackathon.vercel.app     ·     0 overlapping bookings, ever",
     size=12, color=RGBColor(0xB8, 0xB0, 0xA4), caps=True)

# ── 02 · problem ──────────────────────────────────────────────────────────
s = slide(); head(s, "The problem, in the brief's own words",
                  "At 6 PM, fifty students\nwant the same court", "02")
text(s, M, Inches(2.35), Inches(5.5), Inches(2.4),
     [[('"The system must confirm ', {}), ("exactly one", {"bold": True, "color": INK}),
       (' and reject the rest, without corrupting data."', {})],
      [("Almost every booking system answers this the same way — and almost "
        "every one of them is wrong.", {})]],
     size=17, spacing=1.4, space_after=14)
code(s, Inches(6.9), Inches(2.35), Inches(5.7), Inches(1.75),
     [[("// the question", {"color": INK3})],
      [("const taken = await isSlotTaken(court, slot);", {})],
      [("", {})],
      [("// the answer", {"color": INK3})],
      [("if (!taken) await insertBooking(court, slot);", {})]], size=12)
text(s, Inches(6.9), Inches(4.3), Inches(5.7), Inches(1.0),
     [[("Both lines are individually correct. Between them is a ", {}),
       ("gap", {"bold": True, "color": INK}),
       (" — and under load another request commits inside it.", {})]],
     size=15, spacing=1.35)

# ── 03 · the answer ───────────────────────────────────────────────────────
s = slide(); head(s, "The answer", "The write IS the decision", "03")
code(s, M, Inches(2.3), CW, Inches(1.6),
     [[("ALTER TABLE bookings ADD CONSTRAINT bookings_no_overlap", {})],
      [("  EXCLUDE USING gist (", {"color": RUST, "bold": True}),
       ("facility_id ", {}), ("WITH =", {"color": RUST, "bold": True}),
       (", during ", {}), ("WITH &&", {"color": RUST, "bold": True}),
       (")", {"color": RUST, "bold": True})],
      [("  WHERE (status = 'confirmed');", {})]], size=15)
text(s, M, Inches(4.2), CW, Inches(2.0),
     [[("There is no question and no answer. A booking carries a time range, and "
        "the database ", {}),
       ("physically cannot store two overlapping ranges", {"bold": True, "color": INK}),
       (" for the same facility. The loser is rejected with SQLSTATE ", {}),
       ("23P01", {"bold": True, "color": INK, "font": MONO}),
       (" before it ever becomes a row — no availability check on the write path, "
        "and no application lock on the slot.", {})]],
     size=18, spacing=1.45)

# ── 04 · unique vs exclude ────────────────────────────────────────────────
s = slide(); head(s, "Why not a unique key", "Partial overlap is the case that happens", "04")
table(s, M, Inches(2.25), CW,
      [["", "UNIQUE (facility, starts_at)", "EXCLUDE … WITH &&"],
       ["Two bookings at 18:00", ("rejected", {"color": OPEN, "bold": True}),
        ("rejected", {"color": OPEN, "bold": True})],
       [("18:00–19:00 vs 18:30–19:30", {"bold": True, "color": INK}),
        ("both stored", {"color": TAKEN, "bold": True}),
        ("rejected", {"color": OPEN, "bold": True})],
       ["A two-hour closure over four slots", ("not expressible", {"color": TAKEN, "bold": True}),
        ("one row", {"color": OPEN, "bold": True})],
       ["Cancelling frees the slot", "needs a delete", ("status flip", {"color": OPEN, "bold": True})]],
      [0.44, 0.28, 0.28], size=15, row_h=Inches(0.55))
text(s, M, Inches(5.6), CW, Inches(0.9),
     [[("A booking from 18:30 does not share a start time with one from 18:00 — but it ", {}),
       ("overlaps", {"bold": True, "color": INK}),
       (" it, and overlap is what matters to a student holding a racket.", {})]],
     size=17, spacing=1.4)

# ── 05 · infographic (image belongs here) ─────────────────────────────────
s = slide()
s.shapes.add_picture("docs/media/how-it-decides.png", 0, 0, width=W, height=H)

# ── 06 · the two verdicts ─────────────────────────────────────────────────
s = slide(); head(s, "The same burst, fired twice at production",
                  "Fifty requests. One court.", "06", "Measured live · not simulated")
text(s, M, Inches(2.3), Inches(5.6), Inches(0.3), "Check, then write",
     size=13, color=TAKEN, bold=True, caps=True)
stat(s, M, Inches(2.75), "50", "bookings confirmed. Nothing errored.", TAKEN)
stat(s, M, Inches(4.35), "1225", "overlapping pairs, swept across the whole table", TAKEN)
text(s, Inches(7.1), Inches(2.3), Inches(5.5), Inches(0.3), "The write is the decision",
     size=13, color=OPEN, bold=True, caps=True)
stat(s, Inches(7.1), Inches(2.75), "1", "confirmed — 49 rejected with SQLSTATE 23P01", OPEN)
stat(s, Inches(7.1), Inches(4.35), "0", "overlapping pairs, across 380+ confirmed rows", INK)

# ── 07 · measured ─────────────────────────────────────────────────────────
s = slide(); head(s, "Measured on the deployed app", "Sub-linear, not linear", "07",
                  "Twenty times the load, about ten times the time")
text(s, M, Inches(2.3), Inches(4.5), Inches(2.6),
     [[("All contenders attempt the insert at once. One wins; the rest block on "
        "its uncommitted row and are released together the moment it commits. ", {}),
       ("Nothing in the write path serialises on n.", {"bold": True, "color": INK})],
      [("Vercel sin1, Neon ap-southeast-1. Median of three warm runs.", {"size": 14, "color": INK3})]],
     size=16, spacing=1.4, space_after=12)
table(s, Inches(5.7), Inches(2.3), Inches(6.9),
      [["Mode", ("Requests", {"align": "r"}), ("Server", {"align": "r"}),
        ("Confirmed", {"align": "r"}), ("Overlaps", {"align": "r"})],
       ["naive", ("50", {"align": "r", "font": MONO}), ("37 ms", {"align": "r", "font": MONO}),
        ("50", {"align": "r", "font": MONO, "color": TAKEN, "bold": True}),
        ("1225", {"align": "r", "font": MONO, "color": TAKEN, "bold": True})],
       ["safe", ("10", {"align": "r", "font": MONO}), ("47 ms", {"align": "r", "font": MONO}),
        ("1", {"align": "r", "font": MONO}), ("0", {"align": "r", "font": MONO, "color": OPEN, "bold": True})],
       ["safe", ("50", {"align": "r", "font": MONO}), ("116 ms", {"align": "r", "font": MONO}),
        ("1", {"align": "r", "font": MONO}), ("0", {"align": "r", "font": MONO, "color": OPEN, "bold": True})],
       ["safe", ("100", {"align": "r", "font": MONO}), ("251 ms", {"align": "r", "font": MONO}),
        ("1", {"align": "r", "font": MONO}), ("0", {"align": "r", "font": MONO, "color": OPEN, "bold": True})],
       [("safe", {"bold": True, "color": INK}), ("200", {"align": "r", "font": MONO, "bold": True, "color": INK}),
        ("546 ms", {"align": "r", "font": MONO, "bold": True, "color": INK}),
        ("1", {"align": "r", "font": MONO, "bold": True, "color": INK}),
        ("0", {"align": "r", "font": MONO, "color": OPEN, "bold": True})]],
      [0.22, 0.19, 0.19, 0.20, 0.20], size=14, row_h=Inches(0.5))

# ── 08 · beyond correctness ───────────────────────────────────────────────
s = slide(); head(s, 'Beyond "it does not double-book"', "Correctness is the floor", "08")
items = [
    ("Idempotency keys.", " A retried submit returns the original booking. A double-tap on hostel wifi is the common case, not the edge case."),
    ("Typed rejections.", " Losing returns SLOT_TAKEN, QUOTA_EXCEEDED or OVERLAPS_OWN — with three alternative slots and a queue position."),
    ("Waitlist promotion.", " Cancelling frees the slot and offers it to the next student in the same transaction — no window where it is free but unowned."),
    ("Maintenance closures.", " A closure is a booking row. One invariant, two features: a manager cannot close a court out from under a confirmed reservation."),
    ("Fair draw.", " Peak slots go to a seeded, weighted, published lottery. First-come-first-served rewards the best wifi, not the keenest player."),
]
y = Inches(2.3)
for bold, rest in items:
    text(s, M, y, CW, Inches(0.8), [[(bold, {"bold": True, "color": INK}), (rest, {})]],
         size=16, spacing=1.3)
    y += Inches(0.92)

# ── 09 · key features ─────────────────────────────────────────────────────
s = slide(); head(s, "Key features and functionality", "What the product actually does", "09")
table(s, M, Inches(2.2), CW,
      [["Feature", "Functionality"],
       [("Availability view", {"bold": True, "color": INK}), "Twelve facilities with live open, taken and closed counts. Derived from the bookings table per request — never stored, so it cannot drift."],
       [("Slot selection", {"bold": True, "color": INK}), "A day grid per facility on its own timetable, in IST, seven days ahead. Live over Server-Sent Events."],
       [("Booking guarantee", {"bold": True, "color": INK}), "Confirmed by the database or refused by it. Idempotency keys make a retry return the original booking."],
       [("Typed rejections", {"bold": True, "color": INK}), "Each carries a reason, three alternative slots, and a queue position."],
       [("Waitlist", {"bold": True, "color": INK}), "A cancellation offers the slot to the next student in the same transaction, with a claim expiry."],
       [("Fair draw", {"bold": True, "color": INK}), "Peak slots go to a seeded, weighted, published lottery."],
       [("Operations console", {"bold": True, "color": INK}), "Maintenance closures, refused by the same constraint if a student holds that window."],
       [("Insights", {"bold": True, "color": INK}), "Demand by hour, utilisation heatmap, no-show rate, under-used prime slots."]],
      [0.24, 0.76], size=13, row_h=Inches(0.55))

# ── 10 · user flow ────────────────────────────────────────────────────────
s = slide(); head(s, "Overall user flow", "From discovery to confirmation", "10",
                  "Manager flow: close a window → refused if a student holds it → reopening is a status flip")
steps = [("01", "Browse", "Facilities with live open counts. Filter by sport, pick a day."),
         ("02", "Inspect", "The day grid for one court: open, taken, yours, closed."),
         ("03", "Select", "Tap a slot. The sheet shows party size and the idempotency key."),
         ("04", "Submit", "One constrained INSERT. No availability check on the write path."),
         ("05", "Decide", "Postgres confirms exactly one and refuses every other request.")]
bw = (CW - Inches(0.0)) / 5
for i, (n, t, d) in enumerate(steps):
    x = M + bw * i
    box = s.shapes.add_shape(1, x, Inches(2.25), bw, Inches(1.75))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = INK; box.line.width = Pt(1.25); box.shadow.inherit = False
    text(s, x + Inches(0.16), Inches(2.42), bw - Inches(0.32), Inches(0.25), n,
         size=11, color=RUST, bold=True, font=MONO)
    text(s, x + Inches(0.16), Inches(2.68), bw - Inches(0.32), Inches(0.3), t,
         size=15, color=INK, bold=True)
    text(s, x + Inches(0.16), Inches(3.02), bw - Inches(0.32), Inches(0.95), d,
         size=11, color=INK2, spacing=1.2)
for i, (title, body, col) in enumerate([
        ("Won → 201 Confirmed",
         "A booking code is issued. It appears in My Bookings and the grid marks the slot as yours, live on every open screen. Cancelling releases it and promotes the waitlist in one transaction.", OPEN),
        ("Lost → 409 Typed rejection",
         "SQLSTATE 23P01 becomes a reason a student can act on: three alternative slots on the same grid, plus a waitlist place offered automatically when somebody cancels.", TAKEN)]):
    x = M + (CW / 2 + Inches(0.15)) * i
    w = CW / 2 - Inches(0.15)
    bar = s.shapes.add_shape(1, x, Inches(4.35), Pt(4), Inches(1.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    bg = s.shapes.add_shape(1, x + Pt(4), Inches(4.35), w - Pt(4), Inches(1.55))
    bg.fill.solid(); bg.fill.fore_color.rgb = WASH
    bg.line.fill.background(); bg.shadow.inherit = False
    text(s, x + Inches(0.22), Inches(4.5), w - Inches(0.4), Inches(0.3), title,
         size=13, color=col, bold=True, caps=True)
    text(s, x + Inches(0.22), Inches(4.85), w - Inches(0.4), Inches(1.0), body,
         size=13, color=INK2, spacing=1.3)

# ── 11 · wireframes / UI ──────────────────────────────────────────────────
s = slide(); head(s, "Wireframes and UI design", "The screens, as built", "11",
                  "Mobile first — students book from phones")
shots = [("docs/media/home.jpg", "Browse. Hero, live counts, facility grid."),
         ("docs/media/facility.png", "Slot grid. The day as a fixture table."),
         ("docs/media/bookings.png", "My bookings. Upcoming and history."),
         ("docs/media/ops.png", "Ops. Closures, refused on a clash.")]
gw = (CW - Inches(0.45)) / 4
for i, (path, cap) in enumerate(shots):
    x = M + (gw + Inches(0.15)) * i
    s.shapes.add_picture(path, x, Inches(2.25), width=gw)
    text(s, x, Inches(4.05), gw, Inches(0.6), cap, size=11, color=INK2, spacing=1.2)
text(s, M, Inches(4.95), CW / 2 - Inches(0.2), Inches(1.1),
     [[("Design system. ", {"bold": True, "color": INK}),
       ("No brand hue: charcoal chrome, warm paper, rust for the single action, "
        "brass for a hairline. Every saturated colour on screen means something.", {})]],
     size=14, spacing=1.35)
text(s, M + CW / 2 + Inches(0.2), Inches(4.95), CW / 2 - Inches(0.2), Inches(1.1),
     [[("Accessible by construction. ", {"bold": True, "color": INK}),
       ("State is carried by fill rather than hue, and availability is stated three "
        "ways — dot, word, count. Motion is dropped under prefers-reduced-motion.", {})]],
     size=14, spacing=1.35)

# ── 12 · tech stack ───────────────────────────────────────────────────────
s = slide(); head(s, "Technology stack", "What it is built on", "12",
                  "No ORM — the constraint is the product, and it is written in SQL")
app = [["Application", ""],
       [("Next.js 15", {"bold": True, "color": INK}), "App Router, React Server Components"],
       [("React 19", {"bold": True, "color": INK}), "UI runtime"],
       [("TypeScript", {"bold": True, "color": INK}), "Strict mode"],
       [("Tailwind CSS v4", {"bold": True, "color": INK}), "Token layer, no component library"],
       [("zod", {"bold": True, "color": INK}), "Validation on every request body"],
       [("lucide-react / date-fns", {"bold": True, "color": INK}), "Icons; IST-correct date handling"]]
data = [["Data and platform", ""],
        [("PostgreSQL 17", {"bold": True, "color": INK}), "Neon, ap-southeast-1"],
        [("btree_gist", {"bold": True, "color": INK}), "Extension backing the exclusion constraint"],
        [("postgres.js", {"bold": True, "color": INK}), "Driver — tagged templates, bind parameters"],
        [("Vercel", {"bold": True, "color": INK}), "Serverless functions, region sin1"],
        [("Server-Sent Events", {"bold": True, "color": INK}), "Live availability, native"],
        [("Vitest / Puppeteer / CI", {"bold": True, "color": INK}), "19 unit tests, 11 browser checks, real Postgres"]]
table(s, M, Inches(2.2), CW / 2 - Inches(0.3), app, [0.42, 0.58], size=13, row_h=Inches(0.52))
table(s, M + CW / 2 + Inches(0.3), Inches(2.2), CW / 2 - Inches(0.3), data, [0.42, 0.58], size=13, row_h=Inches(0.52))

# ── 13 · the product ──────────────────────────────────────────────────────
s = slide(); head(s, "The product", "Built for a student on a phone at 5:58 PM", "13",
                  "Twelve facilities · availability derived, never stored")
prod = [("docs/media/facility.png", "Booking. The day as a fixture table, live over SSE."),
        ("docs/media/fair.png", "Fair draw. Rules and seed published before the draw."),
        ("docs/media/insights.png", "Insights. Aggregated from the booking table itself.")]
gw = (CW - Inches(0.6)) / 3
for i, (path, cap) in enumerate(prod):
    x = M + (gw + Inches(0.3)) * i
    s.shapes.add_picture(path, x, Inches(2.3), width=gw)
    text(s, x, Inches(4.75), gw, Inches(0.6), cap, size=13, color=INK2, spacing=1.25)

# ── 14 · verification ─────────────────────────────────────────────────────
s = slide(); head(s, "Verified, not asserted", "Every claim is reproducible", "14")
stat(s, M, Inches(2.4), "11/11", "browser journey against production — book, confirm, cancel, slot returns", OPEN, 44)
stat(s, M, Inches(3.9), "19", "unit tests, run in CI against a real Postgres 17", OPEN, 44)
stat(s, M, Inches(5.4), "0", "npm audit vulnerabilities", OPEN, 44)
code(s, Inches(7.1), Inches(2.4), Inches(5.5), Inches(1.5),
     [[("npm run test       ", {}), ("# concurrency + lottery", {"color": INK3})],
      [("npm run invariant  ", {}), ("# whole-table sweep", {"color": INK3})],
      [("npm run e2e        ", {}), ("# 11 browser checks", {"color": INK3})]], size=13)
text(s, Inches(7.1), Inches(4.2), Inches(5.5), Inches(1.4),
     [[("The overlap count is never a check of the slot just tested — it is ", {}),
       ("every confirmed row against every other row", {"bold": True, "color": INK}),
       (" on the same court, run after every race.", {})]], size=15, spacing=1.4)

# ── 15 · scope ────────────────────────────────────────────────────────────
s = slide(); head(s, "Stated plainly", "What is deliberately out of scope", "15")
scope = [
    ("No passwords.", " Identity is a signed cookie over a seeded roster with a switch-reader control, so a judge can race two students without two browser profiles. Swapping in SSO replaces one function."),
    ("The demo endpoints are unauthenticated", " and rate limited rather than gated, because the race demo has to stay one click away."),
    ("What is enforced:", " ownership scoped inside every UPDATE, zod on every request body, bind parameters on every query, no sql.unsafe anywhere."),
    ("Next:", " institute SSO, push notifications for waitlist offers, and a rate limiter keyed on identity rather than on a header the client controls."),
]
y = Inches(2.35)
for bold, rest in scope:
    text(s, M, y, CW, Inches(1.0), [[(bold, {"bold": True, "color": INK}), (rest, {})]],
         size=16, spacing=1.35)
    y += Inches(1.12)

# ── 16 · close ────────────────────────────────────────────────────────────
s = slide()
text(s, M, Inches(2.3), CW, Inches(0.3), "One platform. Zero clashes.",
     size=14, color=RUST, bold=True, caps=True, align=PP_ALIGN.CENTER)
text(s, M, Inches(2.85), CW, Inches(1.6), "Fifty students. One court.\nOne winner.",
     size=46, color=INK, font=SERIF, bold=True, align=PP_ALIGN.CENTER, spacing=1.05)
rule(s, (W - Inches(0.9)) / 2, Inches(4.62))
text(s, M, Inches(4.95), CW, Inches(0.9),
     "Run the race yourself — the burst is fired at a real Postgres table through "
     "a real HTTP endpoint, and the whole table is swept afterwards.",
     size=17, color=INK2, align=PP_ALIGN.CENTER, spacing=1.4)
text(s, M, Inches(5.95), CW, Inches(0.5), "innovait-hackathon.vercel.app/race",
     size=24, color=RUST, font=MONO, bold=True, align=PP_ALIGN.CENTER)
text(s, M, Inches(6.55), CW, Inches(0.3), "Team InnovAIT · PlayHack SDE Track",
     size=12, color=INK3, caps=True, align=PP_ALIGN.CENTER)

out = "docs/PlayHack.pptx"
prs.save(out)
print(f"{out}  ·  {len(prs.slides._sldIdLst)} slides  ·  {os.path.getsize(out)//1024} KB  ·  native text")
