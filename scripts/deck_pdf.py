"""
Build docs/PlayHack.pdf from the rendered slides, and docs/PlayHack-animated.pptx
which adds the two motion slides (the race demo and the insights sweep) whose
GIFs play in PowerPoint's slideshow mode.

  npm run deck:extras      (after npm run deck)
"""
import glob, os
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDES = sorted(glob.glob("docs/deck/slides/slide-*.png"))
if not SLIDES:
    raise SystemExit("no slides — run `npm run deck` first")

# ── 1. PDF ────────────────────────────────────────────────────────────────
# Straight from the rendered PNGs rather than by converting the .pptx: there is
# no PowerPoint or LibreOffice on this machine, and these are the same pixels
# the .pptx carries, so the PDF cannot drift from the deck.
pages = [Image.open(p).convert("RGB") for p in SLIDES]
pages[0].save("docs/PlayHack.pdf", save_all=True, append_images=pages[1:],
              resolution=144.0)
print(f"docs/PlayHack.pdf  ·  {len(pages)} pages  ·  {os.path.getsize('docs/PlayHack.pdf')//1024} KB")

# ── 2. Animated deck ──────────────────────────────────────────────────────
W, H = Emu(12192000), Emu(6858000)
PAPER, INK, INK3, RUST = RGBColor(0xF7,0xF4,0xEE), RGBColor(0x22,0x20,0x1E), RGBColor(0x83,0x7C,0x73), RGBColor(0xB4,0x51,0x2A)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]

def still(png):
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(png, 0, 0, width=W, height=H)
    return s

def motion(gif, kicker, title, caption):
    """A GIF slide. PowerPoint plays the GIF in slideshow mode; in the editor
    it shows the first frame, which is why the caption has to carry the point
    on its own."""
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_textbox(0, 0, W, H)          # paper ground
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background()

    def text(x, y, w, h, body, size, color, bold=False, caps=False, space=0):
        tb = s.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = body.upper() if caps else body
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
        f.name = "Inter"
        return tb

    M = 800100                                       # ~0.87in margin
    text(M, 520000, W - 2*M, 300000, kicker, 12, RUST, bold=True, caps=True)
    t = text(M, 760000, W - 2*M, 700000, title, 30, INK, bold=True)
    t.text_frame.paragraphs[0].runs[0].font.name = "Playfair Display"

    # The GIF, centred in the space that is left.
    with Image.open(gif) as im:
        gw, gh = im.size
    avail_w, avail_h = W - 2*M, Emu(3900000)
    scale = min(avail_w / gw, avail_h / gh)
    dw, dh = int(gw * scale), int(gh * scale)
    s.shapes.add_picture(gif, int((W - dw) / 2), 1700000, width=dw, height=dh)

    text(M, 1700000 + dh + 180000, W - 2*M, 500000, caption, 13, INK3)
    return s

# The proof slides sit next to the race; the product slide next to the charts.
for i, png in enumerate(SLIDES):
    still(png)
    if i == 5:   # after "Fifty requests. One court."
        motion("docs/media/race-demo.gif",
               "Nothing is simulated",
               "The same burst, fired twice at production",
               "Naive first — 50 confirmed, 1225 overlapping pairs. Then the identical burst at the constrained path: exactly one survives. Plays in slideshow mode.")
    if i == 8:   # after "The product"
        motion("docs/media/insights-demo.gif",
               "Insights",
               "Every figure aggregates the bookings table itself",
               "Demand by hour, utilisation by facility and hour, no-show rate, under-used prime slots. No reporting copy and no nightly export to drift out of step.")

out = "docs/PlayHack-animated.pptx"
prs.save(out)
print(f"{out}  ·  {len(prs.slides.__iter__.__self__._sldIdLst)} slides  ·  {os.path.getsize(out)//1024} KB")
